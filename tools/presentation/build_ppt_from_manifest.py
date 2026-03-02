#!/usr/bin/env python3
"""Build a scientific PowerPoint deck from a manifest JSON."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SUPPORTED_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build PPTX from research presentation manifest JSON.")
    parser.add_argument("--manifest", required=True, help="Path to the manifest JSON.")
    parser.add_argument("--output", required=True, help="Output .pptx file path.")
    return parser.parse_args()


def hex_to_rgb(hex_color: str, fallback: Tuple[int, int, int]) -> RGBColor:
    value = (hex_color or "").strip().lstrip("#")
    if len(value) != 6:
        return RGBColor(*fallback)
    try:
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except ValueError:
        return RGBColor(*fallback)


def resolve_figure_path(path_str: Optional[str], scan_root: Path, manifest_dir: Path) -> Optional[Path]:
    if not path_str:
        return None
    candidate = Path(path_str).expanduser()
    if candidate.is_absolute() and candidate.exists():
        return candidate

    from_scan_root = (scan_root / candidate).resolve()
    if from_scan_root.exists():
        return from_scan_root

    from_manifest_dir = (manifest_dir / candidate).resolve()
    if from_manifest_dir.exists():
        return from_manifest_dir

    return None


def style_with_defaults(style: Dict[str, Any]) -> Dict[str, Any]:
    out = {
        "font_name": style.get("font_name", "Arial"),
        "title_font_size_pt": int(style.get("title_font_size_pt", 24)),
        "body_font_size_pt": int(style.get("body_font_size_pt", 18)),
        "bottom_line_font_size_pt": int(style.get("bottom_line_font_size_pt", 20)),
        "title_bar_bg": style.get("title_bar_bg", "000000"),
        "title_text_color": style.get("title_text_color", "FFFFFF"),
        "bottom_bar_bg": style.get("bottom_bar_bg", "000000"),
        "bottom_text_color": style.get("bottom_text_color", "FFFFFF"),
        "left_panel_fraction": float(style.get("left_panel_fraction", 0.20)),
    }
    out["left_panel_fraction"] = max(0.15, min(0.35, out["left_panel_fraction"]))
    out["body_font_size_pt"] = max(18, out["body_font_size_pt"])
    return out


def add_bar(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    bg_rgb: RGBColor,
    text_rgb: RGBColor,
    font_name: str,
    font_size_pt: int,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_rgb
    shape.line.color.rgb = bg_rgb

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.20)
    tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.0
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = True
    run.font.color.rgb = text_rgb


def add_bullets(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    bullets: List[str],
    font_name: str,
    font_size_pt: int,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    bullets = bullets or ["Add key points for this slide."]

    for idx, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = f"• {bullet.strip()}"
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(8)
        paragraph.space_before = Pt(0)
        paragraph.line_spacing = 1.08

        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(max(18, font_size_pt))


def add_right_panel_image_or_placeholder(
    slide,
    figure_path: Optional[Path],
    left: int,
    top: int,
    width: int,
    height: int,
    font_name: str,
    body_font_size_pt: int,
) -> None:
    # Panel border to keep layout consistent even when image is missing.
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    border.fill.background()
    border.line.color.rgb = RGBColor(60, 60, 60)
    border.line.width = Pt(1.5)

    if figure_path and figure_path.exists() and figure_path.suffix.lower() in SUPPORTED_IMAGE_EXTS:
        pic = slide.shapes.add_picture(str(figure_path), left, top)
        scale = min(width / pic.width, height / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(left + (width - pic.width) / 2)
        pic.top = int(top + (height - pic.height) / 2)
        return

    tf = border.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = "Add figure, result plot, or equation graphic"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(max(18, body_font_size_pt))


def create_deck(manifest: Dict[str, Any], output_path: Path, manifest_dir: Path) -> None:
    style = style_with_defaults(manifest.get("style", {}))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    title_h = Inches(0.82)
    bottom_h = Inches(0.78)
    side_margin = Inches(0.32)
    col_gap = Inches(0.18)
    body_top = title_h + Inches(0.08)
    body_bottom = slide_h - bottom_h - Inches(0.06)
    body_h = int(body_bottom - body_top)

    usable_w = int(slide_w - (2 * side_margin))
    left_w = int(usable_w * style["left_panel_fraction"])
    right_w = int(usable_w - left_w - col_gap)

    left_x = int(side_margin)
    right_x = int(left_x + left_w + col_gap)

    scan_root = Path(manifest.get("scan_root", manifest_dir)).expanduser().resolve()

    title_bg = hex_to_rgb(style["title_bar_bg"], (0, 0, 0))
    title_fg = hex_to_rgb(style["title_text_color"], (255, 255, 255))
    bottom_bg = hex_to_rgb(style["bottom_bar_bg"], (0, 0, 0))
    bottom_fg = hex_to_rgb(style["bottom_text_color"], (255, 255, 255))

    slides = manifest.get("slides", [])
    if not slides:
        raise ValueError("Manifest contains no slides")

    for slide_data in slides:
        slide = prs.slides.add_slide(slide_layout)

        title = str(slide_data.get("title", "Untitled"))
        bottom_line = str(slide_data.get("bottom_line", ""))
        bullets = [str(x) for x in slide_data.get("bullets", []) if str(x).strip()]

        add_bar(
            slide=slide,
            left=0,
            top=0,
            width=slide_w,
            height=title_h,
            text=title,
            bg_rgb=title_bg,
            text_rgb=title_fg,
            font_name=style["font_name"],
            font_size_pt=style["title_font_size_pt"],
        )

        add_bullets(
            slide=slide,
            left=left_x,
            top=int(body_top),
            width=left_w,
            height=body_h,
            bullets=bullets,
            font_name=style["font_name"],
            font_size_pt=style["body_font_size_pt"],
        )

        figure = resolve_figure_path(
            path_str=slide_data.get("figure"),
            scan_root=scan_root,
            manifest_dir=manifest_dir,
        )

        add_right_panel_image_or_placeholder(
            slide=slide,
            figure_path=figure,
            left=right_x,
            top=int(body_top),
            width=right_w,
            height=body_h,
            font_name=style["font_name"],
            body_font_size_pt=style["body_font_size_pt"],
        )

        add_bar(
            slide=slide,
            left=0,
            top=int(slide_h - bottom_h),
            width=slide_w,
            height=bottom_h,
            text=f"Bottom line: {bottom_line}",
            bg_rgb=bottom_bg,
            text_rgb=bottom_fg,
            font_name=style["font_name"],
            font_size_pt=style["bottom_line_font_size_pt"],
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    args = parse_args()

    manifest_path = Path(args.manifest).expanduser().resolve()
    if not manifest_path.exists():
        raise SystemExit(f"manifest not found: {manifest_path}")

    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    output_path = Path(args.output).expanduser().resolve()

    create_deck(manifest=manifest, output_path=output_path, manifest_dir=manifest_path.parent)

    print(f"[ok] pptx written: {output_path}")


if __name__ == "__main__":
    main()
